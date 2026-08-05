"""PPTX knowledge handler."""
from __future__ import annotations

import base64
import io
import logging
from typing import Any, Dict, Optional

from skills.knowledge import register_skill

logger = logging.getLogger(__name__)

MIME_TYPES = [
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
]


async def extract(envelope: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    extra = envelope.get("extra", {})
    mimetype = (
        extra.get("doc_mimetype")
        or "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ).lower()
    source_name = extra.get("doc_file_name") or "document.pptx"
    raw_bytes = await _download_bytes(envelope)
    if raw_bytes is None:
        return None
    text = ""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(raw_bytes))
        parts = []
        for idx, slide in enumerate(prs.slides):
            parts.append(f"--- Slide {idx + 1} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            parts.append(t)
                if shape.has_table:
                    parts.append("[TABELA]")
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        parts.append(" | ".join(cells))
        text = "\n".join(parts)
    except Exception as exc:
        logger.warning("pptx_handler extract failed: %s", exc)
        return None
    return {
        "text": text,
        "source_name": source_name,
        "mimetype": mimetype,
        "raw_size": len(raw_bytes),
    }


async def _download_bytes(envelope: Dict[str, Any]) -> Optional[bytes]:
    extra = envelope.get("extra", {})
    doc_b64 = extra.get("doc_base64", "")
    if doc_b64:
        try:
            return base64.b64decode(doc_b64)
        except Exception:
            pass
    message_id = envelope.get("message_id", "")
    instance = envelope.get("instance", "")
    remote_jid = extra.get("remote_jid", "")
    if message_id and instance:
        try:
            from core.evolution_client import get_base64_from_media_message

            result = await get_base64_from_media_message(
                instance=instance,
                message_id=message_id,
                remote_jid=remote_jid,
            )
            if result and result.get("base64"):
                return base64.b64decode(result["base64"])
        except Exception as exc:
            logger.warning("pptx_handler download failed: %s", exc)
    return None


register_skill("pptx", type("PPTXSkill", (), {"MIME_TYPES": MIME_TYPES, "extract": extract}))

__all__ = ["extract", "MIME_TYPES"]
